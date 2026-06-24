from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from insights_generator import generate_insights
from performance_analyzer import analyze_performance

from chart_generator import (
    create_revenue_chart,
    create_region_chart,
    create_product_chart
)


def create_ppt(df):

    # ==========================
    # Extract Sheets
    # ==========================

    employees_df = df["Employees"]
    revenue_df = df["Monthly Revenue"]
    products_df = df["Products"]
    regional_df = df["Regional Performance"]

    # ==========================
    # Analysis
    # ==========================

    performance_df = analyze_performance(
        employees_df
    )

    total_employees = len(employees_df)

    total_sales = employees_df["Sales"].sum()

    top_performer = employees_df.loc[
        employees_df["Sales"].idxmax(),
        "Name"
    ]

    best_region = regional_df.loc[
        regional_df["Revenue"].idxmax(),
        "Region"
    ]

    total_product_revenue = (
        products_df["Revenue"].sum()
    )

    insights = generate_insights(
        employees_df,
        revenue_df,
        products_df,
        regional_df
    )

    # ==========================
    # Generate Charts
    # ==========================

    revenue_chart = create_revenue_chart(
        revenue_df
    )

    region_chart = create_region_chart(
        regional_df
    )

    product_chart = create_product_chart(
        products_df
    )

    # ==========================
    # Create PPT
    # ==========================

    ppt = Presentation()

    # ==================================================
    # Slide 1 - Cover Page
    # ==================================================

    slide = ppt.slides.add_slide(
        ppt.slide_layouts[0]
    )

    slide.shapes.title.text = (
        "Company Performance Report"
    )

    slide.placeholders[1].text = (
        "Generated Automatically using Python"
    )

    # ==================================================
    # Slide 2 - Executive Summary
    # ==================================================

    slide = ppt.slides.add_slide(
        ppt.slide_layouts[1]
    )

    slide.shapes.title.text = (
        "Executive Summary"
    )

    slide.placeholders[1].text = (
        f"Total Employees: {total_employees}\n"
        f"Total Sales: {total_sales:,}\n"
        f"Top Performer: {top_performer}\n"
        f"Best Region: {best_region}\n"
        f"Product Revenue: {total_product_revenue:,}"
    )

    # ==================================================
    # Slide 3 - KPI Dashboard
    # ==================================================

    slide = ppt.slides.add_slide(
        ppt.slide_layouts[5]
    )

    slide.shapes.title.text = (
        "KPI Dashboard"
    )

    kpis = [
        ("Employees", str(total_employees)),
        ("Sales", f"{total_sales:,}"),
        ("Top Performer", top_performer),
        ("Best Region", best_region)
    ]

    positions = [
        (0.5, 1.5),
        (5.0, 1.5),
        (0.5, 3.5),
        (5.0, 3.5)
    ]

    for i in range(4):

        left = Inches(positions[i][0])
        top = Inches(positions[i][1])

        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            top,
            Inches(3),
            Inches(1.5)
        )

        card.text = (
            f"{kpis[i][0]}\n\n{kpis[i][1]}"
        )

    # ==================================================
    # Slide 4 - Performance Analysis
    # ==================================================

    slide = ppt.slides.add_slide(
        ppt.slide_layouts[5]
    )

    slide.shapes.title.text = (
        "Performance Analysis"
    )

    rows = len(performance_df) + 1
    cols = 4

    table = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.5),
        Inches(1.2),
        Inches(8),
        Inches(2.5)
    ).table

    headers = [
        "Employee",
        "Sales",
        "Target",
        "Achievement %"
    ]

    for col in range(cols):
        table.cell(0, col).text = headers[col]

    for i, (_, row) in enumerate(
        performance_df.iterrows(),
        start=1
    ):
        table.cell(i, 0).text = str(row["Name"])
        table.cell(i, 1).text = str(row["Sales"])
        table.cell(i, 2).text = str(row["Target"])
        table.cell(i, 3).text = (
            f"{row['Achievement %']:.1f}%"
        )

    # ==================================================
    # Slide 5 - Revenue Trend
    # ==================================================

    slide = ppt.slides.add_slide(
        ppt.slide_layouts[5]
    )

    slide.shapes.title.text = (
        "Revenue Trend"
    )

    slide.shapes.add_picture(
        revenue_chart,
        Inches(1),
        Inches(1.3),
        width=Inches(7)
    )

    # ==================================================
    # Slide 6 - Region Chart
    # ==================================================

    slide = ppt.slides.add_slide(
        ppt.slide_layouts[5]
    )

    slide.shapes.title.text = (
        "Regional Revenue Distribution"
    )

    slide.shapes.add_picture(
        region_chart,
        Inches(1),
        Inches(1.3),
        width=Inches(6)
    )

    # ==================================================
    # Slide 7 - Product Chart
    # ==================================================

    slide = ppt.slides.add_slide(
        ppt.slide_layouts[5]
    )

    slide.shapes.title.text = (
        "Product Revenue Analysis"
    )

    slide.shapes.add_picture(
        product_chart,
        Inches(1),
        Inches(1.3),
        width=Inches(7)
    )

    # ==================================================
    # Slide 8 - Generated Insights
    # ==================================================

    slide = ppt.slides.add_slide(
        ppt.slide_layouts[1]
    )

    slide.shapes.title.text = (
        "Generated Insights"
    )

    slide.placeholders[1].text = "\n\n".join(
        [f"• {insight}" for insight in insights]
    )

    # ==================================================
    # Employee Detail Tables
    # ==================================================

    for _, row in employees_df.iterrows():

        slide = ppt.slides.add_slide(
            ppt.slide_layouts[5]
        )

        slide.shapes.title.text = (
            f"Employee Profile - {row['Name']}"
        )

        rows = len(employees_df.columns) + 1
        cols = 2

        table = slide.shapes.add_table(
            rows,
            cols,
            Inches(0.5),
            Inches(1.2),
            Inches(8),
            Inches(4)
        ).table

        table.cell(0, 0).text = "Field"
        table.cell(0, 1).text = "Value"

        for i, column in enumerate(
            employees_df.columns,
            start=1
        ):
            table.cell(i, 0).text = str(column)
            table.cell(i, 1).text = str(
                row[column]
            )

    return ppt